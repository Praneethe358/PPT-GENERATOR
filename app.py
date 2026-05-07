"""Auto PPT generator
Random change 1: header added"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

def add_slide(title, points, module_text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # Set background color (dark)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(20, 20, 20)

    # Title
    title_box = slide.shapes.title
    title_box.text = title
    for p in title_box.text_frame.paragraphs:
        for run in p.runs:
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.size = Pt(32)

    # Content
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Header (top-right)
    left = Inches(7)
    top = Inches(0.2)
    width = Inches(3)
    height = Inches(0.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.text = module_text

    for p in tf.paragraphs:
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(200, 200, 200)


# Example input
module_text = "Module 4 | URK24AI1036"

slides_data = [
    ("Introduction to Data Architecture", [
        "Defines how data flows in a system",
        "Helps in organizing data efficiently",
        "Important for scalability",
        "Used in all modern applications"
    ]),
    ("AWS Well-Architected Framework", [
        "Guide for building strong systems",
        "Includes security and reliability",
        "Improves performance",
        "Helps reduce cost"
    ]),
    ("Modern Data Architecture", [
        "Uses cloud-based systems",
        "Supports real-time data",
        "Flexible and scalable",
        "Combines lakes and warehouses"
    ]),
    ("Data Pipeline Design Patterns", [
        "Reusable design approaches",
        "Batch processing",
        "Streaming pipelines",
        "Lambda architecture"
    ]),
    ("Building Scalable Pipelines", [
        "Handles large data",
        "Ensures fault tolerance",
        "Optimizes cost",
        "Maintains performance"
    ]),
    ("Key Takeaways", [
        "Good design is important",
        "Frameworks guide development",
        "Scalability is key",
        "Strong base for future modules"
    ])
]

# Generate slides
for title, points in slides_data:
    add_slide(title, points, module_text)

# Save file
prs.save("Module_4_Presentation.pptx")
print("praneeth")
print("praneeth")
print("praneeth")